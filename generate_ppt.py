from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x1A, 0x25, 0x3C)
ACCENT     = RGBColor(0xFF, 0x6B, 0x35)
ACCENT2    = RGBColor(0xFF, 0x9F, 0x1C)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY  = RGBColor(0xA0, 0xAE, 0xC4)
BLUE       = RGBColor(0x4E, 0xA8, 0xDE)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE     = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_text(slide, left, top, width, height, text, size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def add_bullet_list(slide, left, top, width, height, items, size=16, color=TEXT_WHITE, bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run1 = p.add_run()
        run1.text = "▸ "
        run1.font.size = Pt(size)
        run1.font.color.rgb = bullet_color
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
    return tb

def add_accent_line(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    return s

def slide_header(slide, title, subtitle=None):
    add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(2))
    add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.7), title, size=32, bold=True, color=TEXT_WHITE)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.35), Inches(10), Inches(0.5), subtitle, size=16, color=TEXT_GRAY)

def add_card(slide, left, top, width, height, title, items, icon="", title_color=ACCENT):
    add_shape(slide, left, top, width, height, BG_CARD)
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
             f"{icon}  {title}", size=18, bold=True, color=title_color)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                    height - Inches(0.9), items, size=14, color=TEXT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s1)
# Decorative shapes
add_shape(s1, Inches(-1), Inches(-1), Inches(6), Inches(9.5), RGBColor(0xFF, 0x6B, 0x35))
add_shape(s1, Inches(-1.2), Inches(-1), Inches(6), Inches(9.5), RGBColor(0x14, 0x1E, 0x33))
# Fire emoji
add_text(s1, Inches(1), Inches(1.2), Inches(3), Inches(1.5), "🔥", size=80, align=PP_ALIGN.LEFT)
# Title
add_text(s1, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
         "Sistema de Gestión de Incendios", size=44, bold=True, color=TEXT_WHITE)
add_text(s1, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "Municipalidad Valle del Sol", size=28, color=ACCENT)
add_accent_line(s1, Inches(1), Inches(4.6), Inches(4))
add_text(s1, Inches(1), Inches(5.0), Inches(10), Inches(0.5),
         "Plataforma Fullstack para monitoreo y gestión de incendios forestales", size=16, color=TEXT_GRAY)
add_text(s1, Inches(1), Inches(5.6), Inches(10), Inches(0.5),
         "Ignacio Salazar  ·  Javier Quiroga", size=18, color=TEXT_WHITE, bold=True)
add_text(s1, Inches(1), Inches(6.1), Inches(10), Inches(0.5),
         "DuocUC — Caso Semestral  |  Mayo 2026", size=14, color=TEXT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCCIÓN
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
slide_header(s2, "Introducción", "Presentación del Caso")

add_text(s2, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
         "La Municipalidad Valle del Sol enfrenta el desafío de gestionar y responder eficientemente a los "
         "incendios forestales en su territorio. Se requiere un sistema digital que permita a los ciudadanos "
         "reportar incendios en tiempo real, clasificar automáticamente su gravedad mediante inteligencia "
         "artificial, y coordinar la respuesta de brigadistas desde un panel administrativo centralizado.",
         size=16, color=TEXT_GRAY)

# Problem - Solution cards
add_card(s2, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.2),
         "Problemática", [
             "Reportes manuales lentos e imprecisos",
             "Sin clasificación automática de gravedad",
             "Falta de visibilidad en tiempo real",
             "Coordinación deficiente de brigadistas",
             "Sin registro histórico ni estadísticas"
         ], icon="⚠️", title_color=RED)

add_card(s2, Inches(6.8), Inches(3.5), Inches(5.5), Inches(3.2),
         "Solución Propuesta", [
             "App web con reporte geolocalizado + foto",
             "IA (CNN + XGBoost) clasifica gravedad 1-5",
             "WebSocket para alertas en tiempo real",
             "Dashboard admin con estadísticas",
             "Flujo de estados con auditoría automática"
         ], icon="✅", title_color=GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — SERVICIOS UTILIZADOS
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
slide_header(s3, "Servicios Utilizados", "Plataformas cloud para despliegue y operación")

services = [
    ("☁️", "Vercel", "Frontend Hosting", ["Despliegue automático desde GitHub", "CDN global, SSL gratuito", "Preview deploys por cada PR", "React + Vite optimizado"], BLUE),
    ("🚂", "Railway", "Backend Hosting", ["Contenedor Docker para Spring Boot", "Variables de entorno seguras", "Logs y monitoreo integrado", "Escalamiento automático"], GREEN),
    ("🐘", "Supabase", "Base de Datos", ["PostgreSQL gestionado en la nube", "Extensiones PostGIS habilitadas", "Backups automáticos", "Connection pooling"], PURPLE),
    ("🤗", "HuggingFace Spaces", "Modelo IA", ["Hosting gratuito del modelo CNN", "API REST via Gradio", "GPU inference disponible", "Versionado de modelos"], ACCENT),
]

for i, (icon, name, role, items, color) in enumerate(services):
    left = Inches(0.6 + i * 3.1)
    add_shape(s3, left, Inches(2.2), Inches(2.8), Inches(4.6), BG_CARD)
    add_text(s3, left, Inches(2.3), Inches(2.8), Inches(0.6), icon, size=36, align=PP_ALIGN.CENTER)
    add_text(s3, left, Inches(2.9), Inches(2.8), Inches(0.5), name, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s3, left, Inches(3.4), Inches(2.8), Inches(0.4), role, size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_bullet_list(s3, left + Inches(0.2), Inches(3.9), Inches(2.4), Inches(2.8), items, size=12, color=TEXT_GRAY, bullet_color=color)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — TECNOLOGÍAS UTILIZADAS
# ══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
slide_header(s4, "Tecnologías Utilizadas", "Stack tecnológico del proyecto")

# Frontend card
add_card(s4, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.8),
         "Frontend", [
             "React 19 + TypeScript",
             "Vite 8 (build tool)",
             "React Router DOM 7",
             "Leaflet + React-Leaflet",
             "STOMP.js + SockJS (WebSocket)",
             "Lucide React (iconos)",
             "CSS vanilla (estilos)"
         ], icon="⚛️", title_color=BLUE)

# Backend card
add_card(s4, Inches(4.6), Inches(2.0), Inches(3.8), Inches(4.8),
         "Backend", [
             "Java 21 (LTS)",
             "Spring Boot 3.4.5",
             "Spring Security 6 + JWT",
             "Spring Data JPA + Hibernate",
             "Spring WebSocket (STOMP)",
             "Hibernate Spatial (PostGIS)",
             "Maven 3.9"
         ], icon="☕", title_color=GREEN)

# BD + IA card
add_card(s4, Inches(8.7), Inches(2.0), Inches(3.8), Inches(2.2),
         "Base de Datos", [
             "PostgreSQL 16",
             "PostGIS 3.4 (geoespacial)",
             "uuid-ossp, pg_trgm"
         ], icon="🐘", title_color=PURPLE)

add_card(s4, Inches(8.7), Inches(4.5), Inches(3.8), Inches(2.3),
         "Inteligencia Artificial", [
             "EfficientNet-B0 (CNN/ONNX)",
             "XGBoost (datos satelitales)",
             "Gradio API (serving)"
         ], icon="🧠", title_color=ACCENT)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITECTURA DEL PROYECTO (3 CAPAS)
# ══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
slide_header(s5, "Arquitectura del Proyecto", "Patrón en 3 Capas + Microservicio IA")

# Layer boxes
layers = [
    (Inches(1.5), "🖥️  CAPA DE PRESENTACIÓN", "React + Vite → Vercel",
     ["6 Páginas (Home, Auth, Reportar, Mis Reportes, Detalle, Dashboard)",
      "Componentes: FireMap, Navbar, AlertBanner, Badges",
      "Hooks: useGeolocation, useRealtime (WebSocket STOMP)",
      "Context: AuthContext (sesión JWT persistente)"], BLUE),
    (Inches(3.3), "⚙️  CAPA DE LÓGICA DE NEGOCIO", "Spring Boot → Railway",
     ["Controllers: Auth, Reporte, Dashboard (REST API)",
      "Services: Auth, Reporte, IA, Storage, WebSocket",
      "Security: JWT Filter, Spring Security, BCrypt, Roles",
      "DTOs: Request/Response separados, validación"], GREEN),
    (Inches(5.1), "💾  CAPA DE DATOS", "PostgreSQL → Supabase",
     ["5 Entidades: Usuario, Reporte, Brigadista, HistorialEstado, Alerta",
      "PostGIS: consultas ST_DWithin, ST_MakeEnvelope",
      "Triggers: auto-ubicación, historial de estados",
      "Vistas: estadísticas dashboard, reportes activos"], PURPLE),
]

for top, title, subtitle, items, color in layers:
    add_shape(s5, Inches(0.8), top, Inches(11.5), Inches(1.65), BG_CARD)
    add_text(s5, Inches(1.0), top + Inches(0.08), Inches(5), Inches(0.4), title, size=16, bold=True, color=color)
    add_text(s5, Inches(7), top + Inches(0.08), Inches(5), Inches(0.4), subtitle, size=13, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    add_bullet_list(s5, Inches(1.2), top + Inches(0.45), Inches(10.5), Inches(1.1), items, size=12, color=TEXT_GRAY, bullet_color=color)

# IA side note
add_shape(s5, Inches(9.5), Inches(1.8), Inches(3.0), Inches(1.2), BG_CARD)
add_text(s5, Inches(9.7), Inches(1.85), Inches(2.6), Inches(0.4), "🤗 Microservicio IA", size=14, bold=True, color=ACCENT)
add_text(s5, Inches(9.7), Inches(2.3), Inches(2.6), Inches(0.6), "HuggingFace Space\nGradio API · CNN + XGBoost", size=11, color=TEXT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — ARQUITECTURA DE LOS CÓDIGOS
# ══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)
slide_header(s6, "Arquitectura de los Códigos", "Organización interna del código fuente")

# Backend structure
add_shape(s6, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0), BG_CARD)
add_text(s6, Inches(0.7), Inches(2.1), Inches(5), Inches(0.5), "☕ Backend — Spring Boot", size=20, bold=True, color=GREEN)
code_items = [
    "controller/  → AuthController, ReporteController, DashboardController",
    "service/       → AuthService, ReporteService, IAService, StorageService, WebSocketService",
    "model/          → Usuario, Reporte, Brigadista, HistorialEstado, Alerta + Enums",
    "repository/  → ReporteRepository (consultas geoespaciales JPA)",
    "config/          → SecurityConfig, CorsConfig, WebSocketConfig, JwtUtil",
    "security/      → JwtAuthFilter, UserDetailsServiceImpl",
    "DTO/              → request/ (Login, Register, Reporte) + response/ (Auth, Reporte)",
    "exception/   → GlobalExceptionHandler (400/500 centralizados)"
]
add_bullet_list(s6, Inches(0.7), Inches(2.6), Inches(5.5), Inches(4.2), code_items, size=12, color=TEXT_GRAY, bullet_color=GREEN)

# Frontend structure
add_shape(s6, Inches(6.8), Inches(2.0), Inches(6.0), Inches(5.0), BG_CARD)
add_text(s6, Inches(7.0), Inches(2.1), Inches(5), Inches(0.5), "⚛️ Frontend — React + TypeScript", size=20, bold=True, color=BLUE)
front_items = [
    "pages/             → HomePage, AuthPage, ReportPage, MyReportsPage, ReportDetailPage, DashboardPage",
    "components/  → layout/Navbar, map/FireMap, ui/AlertBanner, ui/Badges",
    "context/          → AuthContext (login, register, logout, sesión JWT)",
    "hooks/             → useGeolocation (GPS), useRealtime (WebSocket STOMP)",
    "lib/                  → api.ts (HTTP client + interceptors), stomp.ts",
    "types/              → index.ts (interfaces TypeScript: Reporte, Usuario, etc.)",
    "App.tsx            → Router + ProtectedRoute + Layout",
    "index.css         → Design system completo (variables CSS, dark theme)"
]
add_bullet_list(s6, Inches(7.0), Inches(2.6), Inches(5.5), Inches(4.2), front_items, size=12, color=TEXT_GRAY, bullet_color=BLUE)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — ESTRATEGIA DE BRANCHING
# ══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7)
slide_header(s7, "Estrategia de Branching", "Flujo de trabajo Git del equipo")

# Main branch
add_shape(s7, Inches(1.5), Inches(2.3), Inches(10), Inches(0.7), RGBColor(0x1B, 0x5E, 0x20))
add_text(s7, Inches(1.7), Inches(2.35), Inches(9), Inches(0.5), "🟢  main  —  Rama de producción (código estable y desplegado)", size=16, bold=True, color=GREEN)

# Review branch
add_shape(s7, Inches(2.5), Inches(3.3), Inches(9), Inches(0.7), RGBColor(0x1A, 0x23, 0x7E))
add_text(s7, Inches(2.7), Inches(3.35), Inches(8), Inches(0.5), "🔵  develop / review  —  Rama de revisión e integración de cambios", size=16, bold=True, color=BLUE)

# Dev branches
add_shape(s7, Inches(3.5), Inches(4.3), Inches(4), Inches(0.7), RGBColor(0x4A, 0x14, 0x8C))
add_text(s7, Inches(3.7), Inches(4.35), Inches(3.5), Inches(0.5), "🟣  feature/ignacio", size=16, bold=True, color=PURPLE)

add_shape(s7, Inches(8), Inches(4.3), Inches(4), Inches(0.7), RGBColor(0x7B, 0x34, 0x1E))
add_text(s7, Inches(8.2), Inches(4.35), Inches(3.5), Inches(0.5), "🟠  feature/javier", size=16, bold=True, color=ACCENT)

# Flow description
add_text(s7, Inches(1.5), Inches(5.3), Inches(10), Inches(0.4), "Flujo de trabajo:", size=18, bold=True, color=TEXT_WHITE)
flow_items = [
    "Cada integrante trabaja en su rama personal (feature/ignacio, feature/javier)",
    "Al completar una funcionalidad, se crea un Pull Request hacia la rama develop/review",
    "Se realiza revisión de código entre compañeros antes de aprobar el merge",
    "Una vez revisado y aprobado, se integra a develop y se testea en conjunto",
    "Cuando develop es estable, se hace merge a main → despliegue automático a producción"
]
add_bullet_list(s7, Inches(1.5), Inches(5.7), Inches(10), Inches(1.6), flow_items, size=14, color=TEXT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — MEJORAS A FUTURO
# ══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8)
slide_header(s8, "Mejoras a Futuro", "Roadmap de evolución del sistema")

# Priority columns
add_card(s8, Inches(0.4), Inches(2.0), Inches(3.9), Inches(5.0),
         "🔴 Prioridad Alta", [
             "Notificaciones Push (FCM)",
             "Upload a Cloud Storage (S3/R2)",
             "Tests unitarios e integración",
             "Refresh Token (JWT)",
             "Validación de imágenes (MIME, tamaño)"
         ], icon="", title_color=RED)

add_card(s8, Inches(4.6), Inches(2.0), Inches(3.9), Inches(5.0),
         "🟡 Prioridad Media", [
             "Asignación automática de brigadistas",
             "Timeline visual de estados",
             "Filtros avanzados en el mapa",
             "Gráficos en Dashboard (Chart.js)",
             "Paginación + Caché (Redis)"
         ], icon="", title_color=ACCENT2)

add_card(s8, Inches(8.8), Inches(2.0), Inches(3.9), Inches(5.0),
         "🟢 Futuro", [
             "PWA / App móvil (React Native)",
             "Integración NASA FIRMS",
             "Modelo IA v2 (dataset chileno)",
             "CI/CD (GitHub Actions)",
             "Monitoreo (Prometheus + Grafana)"
         ], icon="", title_color=GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSIÓN
# ══════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9)
slide_header(s9, "Conclusión")

add_text(s9, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5),
         "Se desarrolló exitosamente un sistema fullstack completo que permite a ciudadanos reportar "
         "incendios forestales con geolocalización y fotografía, clasificar automáticamente la gravedad "
         "mediante inteligencia artificial (CNN + XGBoost), y gestionar la respuesta en tiempo real a "
         "través de un panel administrativo con comunicación WebSocket.",
         size=16, color=TEXT_GRAY)

highlights = [
    ("12", "Servicios\nimplementados"),
    ("3", "Capas de\narquitectura"),
    ("6", "Páginas\nfrontend"),
    ("5", "Niveles de\ngravedad IA"),
]
for i, (num, label) in enumerate(highlights):
    left = Inches(1.0 + i * 2.8)
    add_shape(s9, left, Inches(3.8), Inches(2.3), Inches(1.5), BG_CARD)
    add_text(s9, left, Inches(3.85), Inches(2.3), Inches(0.8), num, size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s9, left, Inches(4.6), Inches(2.3), Inches(0.6), label, size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Deployment summary
add_text(s9, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.4), "Despliegue en producción:", size=18, bold=True, color=TEXT_WHITE)
deploy_items = [
    "Frontend React + Vite → Vercel (CDN global)",
    "Backend Spring Boot Java → Railway (contenedor Docker)",
    "Base de datos PostgreSQL + PostGIS → Supabase (managed)",
    "Modelo IA CNN + XGBoost → HuggingFace Spaces (Gradio API)"
]
add_bullet_list(s9, Inches(0.8), Inches(6.1), Inches(11), Inches(1.2), deploy_items, size=14, color=TEXT_GRAY)

# Thank you
add_text(s9, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
         "Ignacio Salazar  ·  Javier Quiroga  |  DuocUC — Mayo 2026",
         size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# ─── Save ───
output = os.path.join(os.path.dirname(__file__), "Presentacion_Incendios.pptx")
prs.save(output)
print(f"OK - Presentacion guardada en: {output}")
